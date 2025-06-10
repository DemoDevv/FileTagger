import subprocess
from sys import platform

import docx
import pandas as pd
import pptx
import pypdf


def pdf_reader(file_path) -> str:
    with open(file_path, "rb") as file:
        return " ".join(page.extract_text() for page in pypdf.PdfReader(file).pages)


def docx_reader(file_path) -> str:
    return " ".join(
        paragraph.text for paragraph in docx.Document(str(file_path)).paragraphs
    )


def doc_reader(file_path) -> str:
    if platform == "darwin":
        try:
            content = subprocess.run(
                ["textutil", "-stdout", "-cat", "txt", str(file_path)],
                capture_output=True,
                text=True,
            )
            return content.stdout
        except Exception as e:
            print(f"Error reading doc file: {e}")
            return ""
    else:
        raise NotImplementedError("doc_reader is not implemented for this platform")


def xlsx_reader(file_path) -> str:
    return " ".join(pd.read_excel(file_path).astype(str).values.flatten())


def pptx_reader(file_path) -> str:
    prs = pptx.Presentation(str(file_path))

    text_buffer = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                for run in paragraph.runs:
                    text_buffer.append(run.text)

    return " ".join(text_buffer)


def txt_reader(file_path) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()
